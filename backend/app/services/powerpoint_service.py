from __future__ import annotations

from datetime import date, datetime
from io import BytesIO
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy import text

from app.database import engine
from app.services.report_branding_service import (
    ReportBranding,
    get_report_branding,
)


# -------------------------------------------------------------------
# Presentation dimensions
# -------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# -------------------------------------------------------------------
# Fixed supporting colors
# -------------------------------------------------------------------

WHITE = "FFFFFF"
TEXT_DARK = "0F172A"
TEXT_MUTED = "64748B"
BACKGROUND = "F8FAFC"
BORDER = "E2E8F0"

GREEN = "16A34A"
GREEN_LIGHT = "DCFCE7"
AMBER = "D97706"
AMBER_LIGHT = "FEF3C7"
RED = "DC2626"
RED_LIGHT = "FEE2E2"
BLUE = "2563EB"
BLUE_LIGHT = "DBEAFE"


# -------------------------------------------------------------------
# Database helpers
# -------------------------------------------------------------------

def _fetch_rows(query: str) -> List[Dict[str, Any]]:
    """
    Execute a read-only SQL query.

    The PowerPoint export should remain available even when an optional
    operational table has no data or is not yet deployed.
    """

    try:
        with engine.connect() as connection:
            result = connection.execute(text(query))
            return [dict(row._mapping) for row in result]
    except Exception:
        return []


def _fetch_production_data(limit: int = 30) -> List[Dict[str, Any]]:
    return _fetch_rows(
        f"""
        SELECT
            report_date,
            mine_name,
            ore_plan,
            ore_actual,
            waste_plan,
            waste_actual,
            created_at
        FROM public.production_daily
        ORDER BY report_date DESC
        LIMIT {int(limit)}
        """
    )[::-1]


def _fetch_fleet_data(limit: int = 30) -> List[Dict[str, Any]]:
    return _fetch_rows(
        f"""
        SELECT
            report_date,
            mine_name,
            availability,
            utilization,
            created_at
        FROM public.fleet_daily
        ORDER BY report_date DESC
        LIMIT {int(limit)}
        """
    )[::-1]


def _fetch_plant_data(limit: int = 30) -> List[Dict[str, Any]]:
    return _fetch_rows(
        f"""
        SELECT
            report_date,
            mine_name,
            throughput_plan,
            throughput_actual,
            recovery,
            created_at
        FROM public.plant_daily
        ORDER BY report_date DESC
        LIMIT {int(limit)}
        """
    )[::-1]


def _fetch_safety_data(limit: int = 30) -> List[Dict[str, Any]]:
    return _fetch_rows(
        f"""
        SELECT
            report_date,
            mine_name,
            incidents,
            near_misses,
            critical_risks,
            safety_score,
            created_at
        FROM public.safety_daily
        ORDER BY report_date DESC
        LIMIT {int(limit)}
        """
    )[::-1]


def _fetch_executive_actions(limit: int = 8) -> List[Dict[str, Any]]:
    return _fetch_rows(
        f"""
        SELECT
            title,
            priority,
            owner,
            timing,
            expected_benefit,
            status,
            linked_cause
        FROM public.executive_actions
        ORDER BY
            CASE LOWER(COALESCE(priority, ''))
                WHEN 'critical' THEN 1
                WHEN 'high' THEN 2
                WHEN 'medium' THEN 3
                WHEN 'low' THEN 4
                ELSE 5
            END,
            id DESC
        LIMIT {int(limit)}
        """
    )


# -------------------------------------------------------------------
# Value helpers
# -------------------------------------------------------------------

def _number(value: Any) -> float:
    if value is None:
        return 0.0

    try:
        return float(value)
    except (TypeError, ValueError):
        return 0.0


def _safe_ratio(numerator: Any, denominator: Any) -> Optional[float]:
    denominator_value = _number(denominator)

    if denominator_value == 0:
        return None

    return _number(numerator) / denominator_value


def _average(values: Iterable[Any]) -> Optional[float]:
    cleaned = [_number(value) for value in values if value is not None]

    if not cleaned:
        return None

    return sum(cleaned) / len(cleaned)


def _latest(rows: Sequence[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    return rows[-1] if rows else None


def _format_date(value: Any) -> str:
    if isinstance(value, datetime):
        return value.strftime("%d %b")
    if isinstance(value, date):
        return value.strftime("%d %b")
    return str(value or "")


def _format_number(value: Optional[float], decimals: int = 0) -> str:
    if value is None:
        return "No data"

    return f"{value:,.{decimals}f}"


def _format_percent(value: Optional[float], decimals: int = 1) -> str:
    if value is None:
        return "No data"

    return f"{value:.{decimals}f}%"


def _achievement_status(value: Optional[float]) -> Tuple[str, str, str]:
    """
    Return status label, foreground color, and light background color.

    Input is expected as a ratio, for example 0.96 or 1.04.
    """

    if value is None:
        return "No data", TEXT_MUTED, BACKGROUND

    if value >= 1.0:
        return "On / above target", GREEN, GREEN_LIGHT

    if value >= 0.9:
        return "Watch", AMBER, AMBER_LIGHT

    return "Below target", RED, RED_LIGHT


def _hex_to_rgb(value: str) -> RGBColor:
    cleaned = value.strip().replace("#", "")

    if len(cleaned) != 6:
        cleaned = TEXT_DARK

    return RGBColor(
        int(cleaned[0:2], 16),
        int(cleaned[2:4], 16),
        int(cleaned[4:6], 16),
    )


# -------------------------------------------------------------------
# Shape helpers
# -------------------------------------------------------------------

def _set_slide_background(slide, color: str = BACKGROUND) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color)


def _add_text(
    slide,
    text_value: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: float = 18,
    color: str = TEXT_DARK,
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    font_name: str = "Aptos",
) -> Any:
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical_anchor
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment

    run = paragraph.add_run()
    run.text = str(text_value)
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _hex_to_rgb(color)

    return box


def _add_rounded_rectangle(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: str = WHITE,
    line_color: str = BORDER,
    radius_shape: MSO_AUTO_SHAPE_TYPE = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    ),
) -> Any:
    shape = slide.shapes.add_shape(
        radius_shape,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_color)
    shape.line.color.rgb = _hex_to_rgb(line_color)
    shape.line.width = Pt(0.8)

    return shape


def _add_divider(
    slide,
    left: float,
    top: float,
    width: float,
    color: str = BORDER,
) -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.015),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = _hex_to_rgb(color)
    line.line.fill.background()


def _add_slide_number(slide, number: int) -> None:
    _add_text(
        slide,
        str(number),
        12.65,
        7.05,
        0.35,
        0.2,
        font_size=9,
        color=TEXT_MUTED,
        alignment=PP_ALIGN.RIGHT,
    )


def _add_logo(
    slide,
    branding: ReportBranding,
    left: float,
    top: float,
    width: float,
) -> bool:
    if not branding.logo_path:
        return False

    logo_path = Path(branding.logo_path)

    if not logo_path.exists():
        return False

    try:
        slide.shapes.add_picture(
            str(logo_path),
            Inches(left),
            Inches(top),
            width=Inches(width),
        )
        return True
    except Exception:
        return False


def _add_standard_header(
    slide,
    branding: ReportBranding,
    title: str,
    subtitle: Optional[str],
    slide_number: int,
) -> None:
    primary = branding.primary_color_excel

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.12),
        SLIDE_HEIGHT,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _hex_to_rgb(primary)
    accent.line.fill.background()

    _add_text(
        slide,
        branding.company_name,
        0.55,
        0.28,
        7.5,
        0.3,
        font_size=11,
        color=primary,
        bold=True,
    )

    _add_text(
        slide,
        title,
        0.55,
        0.68,
        10.7,
        0.55,
        font_size=27,
        color=TEXT_DARK,
        bold=True,
    )

    if subtitle:
        _add_text(
            slide,
            subtitle,
            0.55,
            1.25,
            10.7,
            0.35,
            font_size=11,
            color=TEXT_MUTED,
        )

    logo_added = _add_logo(
        slide,
        branding,
        left=11.55,
        top=0.25,
        width=1.15,
    )

    if not logo_added:
        _add_text(
            slide,
            branding.mine_name,
            10.1,
            0.35,
            2.55,
            0.35,
            font_size=10,
            color=TEXT_MUTED,
            alignment=PP_ALIGN.RIGHT,
        )

    _add_divider(slide, 0.55, 1.72, 12.1)
    _add_slide_number(slide, slide_number)


def _add_kpi_card(
    slide,
    title: str,
    value: str,
    status: str,
    left: float,
    top: float,
    width: float,
    height: float,
    status_color: str,
    status_background: str,
    note: Optional[str] = None,
) -> None:
    _add_rounded_rectangle(
        slide,
        left,
        top,
        width,
        height,
        fill_color=WHITE,
    )

    _add_text(
        slide,
        title,
        left + 0.25,
        top + 0.18,
        width - 0.5,
        0.35,
        font_size=10,
        color=TEXT_MUTED,
        bold=True,
    )

    _add_text(
        slide,
        value,
        left + 0.25,
        top + 0.64,
        width - 0.5,
        0.55,
        font_size=25,
        color=TEXT_DARK,
        bold=True,
    )

    badge_width = min(max(1.1, len(status) * 0.075), width - 0.5)

    _add_rounded_rectangle(
        slide,
        left + 0.25,
        top + 1.37,
        badge_width,
        0.34,
        fill_color=status_background,
        line_color=status_background,
    )

    _add_text(
        slide,
        status,
        left + 0.34,
        top + 1.405,
        badge_width - 0.18,
        0.24,
        font_size=8.5,
        color=status_color,
        bold=True,
    )

    if note:
        _add_text(
            slide,
            note,
            left + 0.25,
            top + height - 0.48,
            width - 0.5,
            0.28,
            font_size=8.5,
            color=TEXT_MUTED,
        )


def _add_bullet_list(
    slide,
    items: Sequence[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: float = 14,
    color: str = TEXT_DARK,
    bullet_color: Optional[str] = None,
) -> None:
    box = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    for index, item in enumerate(items):
        paragraph = (
            frame.paragraphs[0]
            if index == 0
            else frame.add_paragraph()
        )
        paragraph.text = f"•  {item}"
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = _hex_to_rgb(
            bullet_color or color
        )
        paragraph.space_after = Pt(9)
        paragraph.level = 0


# -------------------------------------------------------------------
# Chart helpers
# -------------------------------------------------------------------

def _add_line_chart(
    slide,
    categories: Sequence[str],
    series: Sequence[Tuple[str, Sequence[float], str]],
    left: float,
    top: float,
    width: float,
    height: float,
    title: Optional[str] = None,
    percentage_axis: bool = False,
) -> None:
    if not categories or not series:
        _add_rounded_rectangle(
            slide,
            left,
            top,
            width,
            height,
            fill_color=WHITE,
        )
        _add_text(
            slide,
            "No trend data available",
            left + 0.3,
            top + height / 2 - 0.2,
            width - 0.6,
            0.4,
            font_size=13,
            color=TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )
        return

    chart_data = CategoryChartData()
    chart_data.categories = list(categories)

    for name, values, _ in series:
        chart_data.add_series(name, list(values))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    ).chart

    chart.has_legend = len(series) > 1

    if chart.has_legend and chart.legend is not None:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)

    chart.has_title = bool(title)

    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = (
        _hex_to_rgb(BORDER)
    )
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)

    if percentage_axis:
        chart.value_axis.tick_labels.number_format = "0%"
        chart.value_axis.maximum_scale = 1.2
        chart.value_axis.minimum_scale = 0

    plot = chart.plots[0]

    for index, (_, _, color) in enumerate(series):
        plot.series[index].format.line.color.rgb = _hex_to_rgb(color)
        plot.series[index].format.line.width = Pt(2)
        plot.series[index].marker.format.fill.solid()
        plot.series[index].marker.format.fill.fore_color.rgb = (
            _hex_to_rgb(color)
        )
        plot.series[index].marker.format.line.color.rgb = (
            _hex_to_rgb(color)
        )


# -------------------------------------------------------------------
# Slide builders
# -------------------------------------------------------------------

def _build_cover_slide(
    presentation: Presentation,
    branding: ReportBranding,
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide, branding.secondary_color_excel)

    primary = branding.primary_color_excel

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.24),
        SLIDE_HEIGHT,
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _hex_to_rgb(primary)
    accent.line.fill.background()

    _add_text(
        slide,
        branding.company_name,
        0.85,
        0.7,
        8.5,
        0.45,
        font_size=16,
        color=primary,
        bold=True,
    )

    _add_text(
        slide,
        "Executive Operations\nBoard Pack",
        0.85,
        1.65,
        8.7,
        1.75,
        font_size=36,
        color=WHITE,
        bold=True,
        vertical_anchor=MSO_ANCHOR.TOP,
    )

    _add_text(
        slide,
        branding.mine_name,
        0.85,
        3.72,
        8.5,
        0.5,
        font_size=18,
        color=WHITE,
    )

    _add_divider(
        slide,
        0.85,
        4.55,
        4.8,
        color=primary,
    )

    generated = datetime.now().strftime("%d %B %Y, %H:%M")

    _add_text(
        slide,
        f"Generated: {generated}\nTimezone: {branding.timezone}",
        0.85,
        4.85,
        5.2,
        0.85,
        font_size=11,
        color="CBD5E1",
        vertical_anchor=MSO_ANCHOR.TOP,
    )

    logo_added = _add_logo(
        slide,
        branding,
        left=10.1,
        top=0.85,
        width=2.1,
    )

    if not logo_added:
        _add_rounded_rectangle(
            slide,
            9.85,
            0.85,
            2.3,
            1.2,
            fill_color=primary,
            line_color=primary,
        )
        _add_text(
            slide,
            branding.company_name,
            10.05,
            1.08,
            1.9,
            0.65,
            font_size=14,
            color=WHITE,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

    _add_text(
        slide,
        "Mine Manager AI",
        10.0,
        6.75,
        2.15,
        0.25,
        font_size=9,
        color="94A3B8",
        alignment=PP_ALIGN.RIGHT,
    )


def _build_kpi_summary_slide(
    presentation: Presentation,
    branding: ReportBranding,
    production_rows: List[Dict[str, Any]],
    fleet_rows: List[Dict[str, Any]],
    plant_rows: List[Dict[str, Any]],
    safety_rows: List[Dict[str, Any]],
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)

    _add_standard_header(
        slide,
        branding,
        "Executive KPI Summary",
        "Current operational position across production, fleet, plant, and safety",
        2,
    )

    latest_production = _latest(production_rows)
    latest_fleet = _latest(fleet_rows)
    latest_plant = _latest(plant_rows)
    latest_safety = _latest(safety_rows)

    ore_achievement = (
        _safe_ratio(
            latest_production.get("ore_actual"),
            latest_production.get("ore_plan"),
        )
        if latest_production
        else None
    )

    fleet_availability = (
        _number(latest_fleet.get("availability"))
        if latest_fleet
        else None
    )

    fleet_utilization = (
        _number(latest_fleet.get("utilization"))
        if latest_fleet
        else None
    )

    plant_achievement = (
        _safe_ratio(
            latest_plant.get("throughput_actual"),
            latest_plant.get("throughput_plan"),
        )
        if latest_plant
        else None
    )

    safety_score = (
        _number(latest_safety.get("safety_score"))
        if latest_safety
        else None
    )

    incidents = (
        int(_number(latest_safety.get("incidents")))
        if latest_safety
        else None
    )

    card_width = 3.82
    card_height = 2.05
    positions = [
        (0.55, 2.05),
        (4.52, 2.05),
        (8.49, 2.05),
        (0.55, 4.38),
        (4.52, 4.38),
        (8.49, 4.38),
    ]

    ore_status = _achievement_status(ore_achievement)
    _add_kpi_card(
        slide,
        "Ore production achievement",
        _format_percent(
            ore_achievement * 100
            if ore_achievement is not None
            else None
        ),
        ore_status[0],
        *positions[0],
        card_width,
        card_height,
        ore_status[1],
        ore_status[2],
        note="Actual ore movement compared with plan",
    )

    availability_ratio = (
        fleet_availability / 100
        if fleet_availability is not None
        else None
    )
    availability_status = _achievement_status(availability_ratio)
    _add_kpi_card(
        slide,
        "Fleet availability",
        _format_percent(fleet_availability),
        availability_status[0],
        *positions[1],
        card_width,
        card_height,
        availability_status[1],
        availability_status[2],
        note="Latest available fleet record",
    )

    utilization_ratio = (
        fleet_utilization / 100
        if fleet_utilization is not None
        else None
    )
    utilization_status = _achievement_status(utilization_ratio)
    _add_kpi_card(
        slide,
        "Fleet utilization",
        _format_percent(fleet_utilization),
        utilization_status[0],
        *positions[2],
        card_width,
        card_height,
        utilization_status[1],
        utilization_status[2],
        note="Latest available fleet record",
    )

    plant_status = _achievement_status(plant_achievement)
    _add_kpi_card(
        slide,
        "Plant throughput achievement",
        _format_percent(
            plant_achievement * 100
            if plant_achievement is not None
            else None
        ),
        plant_status[0],
        *positions[3],
        card_width,
        card_height,
        plant_status[1],
        plant_status[2],
        note="Actual throughput compared with plan",
    )

    safety_ratio = (
        safety_score / 100
        if safety_score is not None
        else None
    )
    safety_status = _achievement_status(safety_ratio)
    _add_kpi_card(
        slide,
        "Safety score",
        _format_percent(safety_score),
        safety_status[0],
        *positions[4],
        card_width,
        card_height,
        safety_status[1],
        safety_status[2],
        note="Latest composite safety score",
    )

    incident_status = (
        ("No incidents", GREEN, GREEN_LIGHT)
        if incidents == 0
        else (
            ("Review required", RED, RED_LIGHT)
            if incidents is not None
            else ("No data", TEXT_MUTED, BACKGROUND)
        )
    )
    _add_kpi_card(
        slide,
        "Safety incidents",
        str(incidents) if incidents is not None else "No data",
        incident_status[0],
        *positions[5],
        card_width,
        card_height,
        incident_status[1],
        incident_status[2],
        note="Latest reporting date",
    )


def _build_production_slide(
    presentation: Presentation,
    branding: ReportBranding,
    production_rows: List[Dict[str, Any]],
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)

    _add_standard_header(
        slide,
        branding,
        "Production Trend",
        "Ore production plan versus actual performance",
        3,
    )

    trend_rows = production_rows[-14:]
    categories = [_format_date(row.get("report_date")) for row in trend_rows]

    plan_values = [_number(row.get("ore_plan")) for row in trend_rows]
    actual_values = [_number(row.get("ore_actual")) for row in trend_rows]

    _add_line_chart(
        slide,
        categories,
        [
            ("Ore Plan", plan_values, TEXT_MUTED),
            (
                "Ore Actual",
                actual_values,
                branding.primary_color_excel,
            ),
        ],
        0.55,
        2.0,
        8.35,
        4.6,
        title="Last 14 reporting periods",
    )

    latest_production = _latest(production_rows)

    if latest_production:
        ore_plan = _number(latest_production.get("ore_plan"))
        ore_actual = _number(latest_production.get("ore_actual"))
        achievement = _safe_ratio(ore_actual, ore_plan)
        variance = ore_actual - ore_plan
    else:
        ore_plan = None
        ore_actual = None
        achievement = None
        variance = None

    status = _achievement_status(achievement)

    _add_rounded_rectangle(
        slide,
        9.2,
        2.0,
        3.55,
        4.6,
        fill_color=WHITE,
    )

    _add_text(
        slide,
        "Latest position",
        9.5,
        2.28,
        2.9,
        0.35,
        font_size=12,
        color=TEXT_DARK,
        bold=True,
    )

    items = [
        ("Ore plan", _format_number(ore_plan)),
        ("Ore actual", _format_number(ore_actual)),
        ("Variance", _format_number(variance)),
        (
            "Achievement",
            _format_percent(
                achievement * 100
                if achievement is not None
                else None
            ),
        ),
    ]

    for index, (label, value) in enumerate(items):
        row_top = 2.9 + index * 0.68

        _add_text(
            slide,
            label,
            9.5,
            row_top,
            1.45,
            0.3,
            font_size=10,
            color=TEXT_MUTED,
        )
        _add_text(
            slide,
            value,
            10.85,
            row_top,
            1.55,
            0.3,
            font_size=13,
            color=TEXT_DARK,
            bold=True,
            alignment=PP_ALIGN.RIGHT,
        )
        _add_divider(
            slide,
            9.5,
            row_top + 0.42,
            2.9,
        )

    _add_rounded_rectangle(
        slide,
        9.5,
        5.72,
        2.9,
        0.5,
        fill_color=status[2],
        line_color=status[2],
    )
    _add_text(
        slide,
        status[0],
        9.7,
        5.82,
        2.5,
        0.25,
        font_size=10,
        color=status[1],
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )


def _build_operations_overview_slide(
    presentation: Presentation,
    branding: ReportBranding,
    fleet_rows: List[Dict[str, Any]],
    plant_rows: List[Dict[str, Any]],
    safety_rows: List[Dict[str, Any]],
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)

    _add_standard_header(
        slide,
        branding,
        "Fleet, Plant and Safety Overview",
        "Recent operational performance across supporting value streams",
        4,
    )

    fleet_trend = fleet_rows[-10:]
    plant_trend = plant_rows[-10:]
    safety_trend = safety_rows[-10:]

    fleet_categories = [
        _format_date(row.get("report_date"))
        for row in fleet_trend
    ]
    fleet_availability = [
        _number(row.get("availability")) / 100
        for row in fleet_trend
    ]
    fleet_utilization = [
        _number(row.get("utilization")) / 100
        for row in fleet_trend
    ]

    _add_line_chart(
        slide,
        fleet_categories,
        [
            ("Availability", fleet_availability, BLUE),
            (
                "Utilization",
                fleet_utilization,
                branding.primary_color_excel,
            ),
        ],
        0.55,
        2.02,
        5.95,
        2.1,
        title="Fleet performance",
        percentage_axis=True,
    )

    plant_categories = [
        _format_date(row.get("report_date"))
        for row in plant_trend
    ]
    plant_achievement = [
        _safe_ratio(
            row.get("throughput_actual"),
            row.get("throughput_plan"),
        )
        or 0
        for row in plant_trend
    ]

    _add_line_chart(
        slide,
        plant_categories,
        [
            (
                "Throughput achievement",
                plant_achievement,
                GREEN,
            )
        ],
        6.8,
        2.02,
        5.95,
        2.1,
        title="Plant throughput achievement",
        percentage_axis=True,
    )

    _add_rounded_rectangle(
        slide,
        0.55,
        4.43,
        12.2,
        1.95,
        fill_color=WHITE,
    )

    latest_safety = _latest(safety_rows)

    safety_metrics = [
        (
            "Incidents",
            str(int(_number(latest_safety.get("incidents"))))
            if latest_safety
            else "No data",
        ),
        (
            "Near misses",
            str(int(_number(latest_safety.get("near_misses"))))
            if latest_safety
            else "No data",
        ),
        (
            "Critical risks",
            str(int(_number(latest_safety.get("critical_risks"))))
            if latest_safety
            else "No data",
        ),
        (
            "Safety score",
            _format_percent(
                _number(latest_safety.get("safety_score"))
                if latest_safety
                else None
            ),
        ),
    ]

    _add_text(
        slide,
        "Latest safety position",
        0.85,
        4.7,
        2.5,
        0.35,
        font_size=12,
        color=TEXT_DARK,
        bold=True,
    )

    for index, (label, value) in enumerate(safety_metrics):
        left = 0.85 + index * 2.9

        _add_text(
            slide,
            label,
            left,
            5.25,
            2.2,
            0.3,
            font_size=9.5,
            color=TEXT_MUTED,
        )
        _add_text(
            slide,
            value,
            left,
            5.62,
            2.2,
            0.42,
            font_size=22,
            color=TEXT_DARK,
            bold=True,
        )


def _build_risk_slide(
    presentation: Presentation,
    branding: ReportBranding,
    production_rows: List[Dict[str, Any]],
    fleet_rows: List[Dict[str, Any]],
    plant_rows: List[Dict[str, Any]],
    safety_rows: List[Dict[str, Any]],
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)

    _add_standard_header(
        slide,
        branding,
        "Key Operational Risks",
        "Automated observations from the latest available operational records",
        5,
    )

    latest_production = _latest(production_rows)
    latest_fleet = _latest(fleet_rows)
    latest_plant = _latest(plant_rows)
    latest_safety = _latest(safety_rows)

    risks: List[Tuple[str, str, str]] = []

    if latest_production:
        ore_achievement = _safe_ratio(
            latest_production.get("ore_actual"),
            latest_production.get("ore_plan"),
        )

        if ore_achievement is not None and ore_achievement < 0.9:
            risks.append(
                (
                    "High",
                    "Ore production materially below plan",
                    (
                        f"Latest achievement is "
                        f"{ore_achievement * 100:.1f}%."
                    ),
                )
            )
        elif ore_achievement is not None and ore_achievement < 1.0:
            risks.append(
                (
                    "Medium",
                    "Ore production below plan",
                    (
                        f"Latest achievement is "
                        f"{ore_achievement * 100:.1f}%."
                    ),
                )
            )

    if latest_fleet:
        availability = _number(latest_fleet.get("availability"))
        utilization = _number(latest_fleet.get("utilization"))

        if availability < 80:
            risks.append(
                (
                    "High",
                    "Fleet availability constraint",
                    f"Latest availability is {availability:.1f}%.",
                )
            )

        if utilization < 75:
            risks.append(
                (
                    "Medium",
                    "Fleet utilization opportunity",
                    f"Latest utilization is {utilization:.1f}%.",
                )
            )

    if latest_plant:
        plant_achievement = _safe_ratio(
            latest_plant.get("throughput_actual"),
            latest_plant.get("throughput_plan"),
        )

        if plant_achievement is not None and plant_achievement < 0.9:
            risks.append(
                (
                    "High",
                    "Plant throughput materially below plan",
                    (
                        f"Latest achievement is "
                        f"{plant_achievement * 100:.1f}%."
                    ),
                )
            )
        elif (
            plant_achievement is not None
            and plant_achievement < 1.0
        ):
            risks.append(
                (
                    "Medium",
                    "Plant throughput below plan",
                    (
                        f"Latest achievement is "
                        f"{plant_achievement * 100:.1f}%."
                    ),
                )
            )

    if latest_safety:
        incidents = int(_number(latest_safety.get("incidents")))
        critical_risks = int(
            _number(latest_safety.get("critical_risks"))
        )

        if incidents > 0:
            risks.append(
                (
                    "Critical",
                    "Safety incident recorded",
                    f"{incidents} incident(s) in the latest record.",
                )
            )

        if critical_risks > 0:
            risks.append(
                (
                    "High",
                    "Open critical safety risks",
                    (
                        f"{critical_risks} critical risk exposure(s) "
                        "require review."
                    ),
                )
            )

    if not risks:
        risks.append(
            (
                "Low",
                "No material threshold exception detected",
                (
                    "Latest records are within the basic automated "
                    "screening thresholds."
                ),
            )
        )

    risks = risks[:5]

    color_map = {
        "Critical": (RED, RED_LIGHT),
        "High": (RED, RED_LIGHT),
        "Medium": (AMBER, AMBER_LIGHT),
        "Low": (GREEN, GREEN_LIGHT),
    }

    for index, (level, title, description) in enumerate(risks):
        top = 2.05 + index * 0.92
        level_color, level_background = color_map[level]

        _add_rounded_rectangle(
            slide,
            0.55,
            top,
            12.2,
            0.72,
            fill_color=WHITE,
        )

        _add_rounded_rectangle(
            slide,
            0.78,
            top + 0.17,
            1.05,
            0.36,
            fill_color=level_background,
            line_color=level_background,
        )
        _add_text(
            slide,
            level,
            0.88,
            top + 0.22,
            0.85,
            0.23,
            font_size=8.5,
            color=level_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        _add_text(
            slide,
            title,
            2.05,
            top + 0.12,
            4.55,
            0.28,
            font_size=11,
            color=TEXT_DARK,
            bold=True,
        )

        _add_text(
            slide,
            description,
            6.45,
            top + 0.12,
            5.95,
            0.4,
            font_size=9.5,
            color=TEXT_MUTED,
        )


def _build_actions_slide(
    presentation: Presentation,
    branding: ReportBranding,
    actions: List[Dict[str, Any]],
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)

    _add_standard_header(
        slide,
        branding,
        "Management Action Register",
        "Priority actions requiring executive attention and ownership",
        6,
    )

    visible_actions = actions[:5]

    if not visible_actions:
        _add_rounded_rectangle(
            slide,
            0.55,
            2.1,
            12.2,
            3.9,
            fill_color=WHITE,
        )
        _add_text(
            slide,
            "No executive actions are currently available.",
            1.0,
            3.4,
            11.3,
            0.5,
            font_size=18,
            color=TEXT_MUTED,
            alignment=PP_ALIGN.CENTER,
        )
        return

    headers = [
        ("Priority", 0.7),
        ("Action", 1.85),
        ("Owner", 7.7),
        ("Timing", 9.5),
        ("Status", 11.05),
    ]

    for label, left in headers:
        _add_text(
            slide,
            label,
            left,
            2.0,
            1.3,
            0.3,
            font_size=9.5,
            color=TEXT_MUTED,
            bold=True,
        )

    for index, action in enumerate(visible_actions):
        top = 2.46 + index * 0.82
        priority = str(action.get("priority") or "Medium").title()
        status = str(action.get("status") or "Open").title()

        priority_colors = {
            "Critical": (RED, RED_LIGHT),
            "High": (RED, RED_LIGHT),
            "Medium": (AMBER, AMBER_LIGHT),
            "Low": (GREEN, GREEN_LIGHT),
        }
        priority_color, priority_background = priority_colors.get(
            priority,
            (TEXT_MUTED, BACKGROUND),
        )

        _add_rounded_rectangle(
            slide,
            0.55,
            top,
            12.2,
            0.64,
            fill_color=WHITE,
        )

        _add_rounded_rectangle(
            slide,
            0.72,
            top + 0.15,
            0.95,
            0.34,
            fill_color=priority_background,
            line_color=priority_background,
        )

        _add_text(
            slide,
            priority,
            0.79,
            top + 0.205,
            0.8,
            0.21,
            font_size=8,
            color=priority_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        _add_text(
            slide,
            str(action.get("title") or "Untitled action"),
            1.85,
            top + 0.12,
            5.55,
            0.42,
            font_size=10,
            color=TEXT_DARK,
            bold=True,
        )

        _add_text(
            slide,
            str(action.get("owner") or "Unassigned"),
            7.7,
            top + 0.12,
            1.55,
            0.4,
            font_size=9,
            color=TEXT_MUTED,
        )

        _add_text(
            slide,
            str(action.get("timing") or "Not set"),
            9.5,
            top + 0.12,
            1.3,
            0.4,
            font_size=9,
            color=TEXT_MUTED,
        )

        _add_text(
            slide,
            status,
            11.05,
            top + 0.12,
            1.35,
            0.4,
            font_size=9,
            color=TEXT_DARK,
            bold=True,
        )


def _build_recommendation_slide(
    presentation: Presentation,
    branding: ReportBranding,
    production_rows: List[Dict[str, Any]],
    fleet_rows: List[Dict[str, Any]],
    plant_rows: List[Dict[str, Any]],
    safety_rows: List[Dict[str, Any]],
) -> None:
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )
    _set_slide_background(slide)

    _add_standard_header(
        slide,
        branding,
        "Executive Recommendations",
        "Recommended management focus for the next operating cycle",
        7,
    )

    recommendations: List[str] = []

    latest_production = _latest(production_rows)
    latest_fleet = _latest(fleet_rows)
    latest_plant = _latest(plant_rows)
    latest_safety = _latest(safety_rows)

    if latest_safety:
        incidents = int(_number(latest_safety.get("incidents")))
        critical_risks = int(
            _number(latest_safety.get("critical_risks"))
        )

        if incidents > 0 or critical_risks > 0:
            recommendations.append(
                "Confirm immediate controls, ownership, and close-out timing "
                "for all current safety incidents and critical risk exposures."
            )

    if latest_production:
        ore_achievement = _safe_ratio(
            latest_production.get("ore_actual"),
            latest_production.get("ore_plan"),
        )

        if ore_achievement is not None and ore_achievement < 1:
            recommendations.append(
                "Review the primary production constraint and agree a "
                "recoverable plan for ore movement before the next shift cycle."
            )

    if latest_fleet:
        availability = _number(latest_fleet.get("availability"))
        utilization = _number(latest_fleet.get("utilization"))

        if availability < 85:
            recommendations.append(
                "Prioritize the equipment availability loss tree and confirm "
                "maintenance recovery actions for the highest-impact assets."
            )

        if utilization < 80:
            recommendations.append(
                "Review dispatch, delay, and operating-time losses to convert "
                "available fleet capacity into productive utilization."
            )

    if latest_plant:
        plant_achievement = _safe_ratio(
            latest_plant.get("throughput_actual"),
            latest_plant.get("throughput_plan"),
        )

        if plant_achievement is not None and plant_achievement < 1:
            recommendations.append(
                "Validate the plant throughput constraint and align mine-to-"
                "mill priorities with the next achievable production target."
            )

    default_recommendations = [
        (
            "Keep one accountable owner and one completion date for each "
            "executive action."
        ),
        (
            "Review exceptions against plan daily and escalate only the "
            "highest-value operational decisions."
        ),
        (
            "Confirm data quality before using KPI movements for management "
            "decisions."
        ),
    ]

    for recommendation in default_recommendations:
        if len(recommendations) >= 5:
            break
        recommendations.append(recommendation)

    _add_rounded_rectangle(
        slide,
        0.55,
        2.05,
        8.2,
        4.65,
        fill_color=WHITE,
    )

    _add_text(
        slide,
        "Recommended next actions",
        0.9,
        2.35,
        4.8,
        0.4,
        font_size=15,
        color=TEXT_DARK,
        bold=True,
    )

    _add_bullet_list(
        slide,
        recommendations[:5],
        0.9,
        3.0,
        7.45,
        3.1,
        font_size=13,
        color=TEXT_DARK,
    )

    primary = branding.primary_color_excel

    _add_rounded_rectangle(
        slide,
        9.05,
        2.05,
        3.7,
        4.65,
        fill_color=branding.secondary_color_excel,
        line_color=branding.secondary_color_excel,
    )

    _add_text(
        slide,
        "Executive focus",
        9.45,
        2.45,
        2.9,
        0.35,
        font_size=13,
        color=primary,
        bold=True,
    )

    _add_text(
        slide,
        "Safety first.\nProtect the plan.\nClose the actions.",
        9.45,
        3.25,
        2.85,
        1.8,
        font_size=24,
        color=WHITE,
        bold=True,
        vertical_anchor=MSO_ANCHOR.TOP,
    )

    _add_text(
        slide,
        (
            "This board pack summarizes current operational records and "
            "supports—not replaces—management judgment."
        ),
        9.45,
        5.75,
        2.85,
        0.55,
        font_size=9,
        color="CBD5E1",
        vertical_anchor=MSO_ANCHOR.TOP,
    )


# -------------------------------------------------------------------
# Public service
# -------------------------------------------------------------------

def generate_executive_powerpoint() -> BytesIO:
    """
    Generate a branded Executive Operations Board Pack.

    Slides:
        1. Cover
        2. Executive KPI Summary
        3. Production Trend
        4. Fleet, Plant and Safety Overview
        5. Key Operational Risks
        6. Management Action Register
        7. Executive Recommendations
    """

    branding = get_report_branding()

    production_rows = _fetch_production_data()
    fleet_rows = _fetch_fleet_data()
    plant_rows = _fetch_plant_data()
    safety_rows = _fetch_safety_data()
    executive_actions = _fetch_executive_actions()

    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    presentation.core_properties.title = (
        f"{branding.company_name} Executive Operations Board Pack"
    )
    presentation.core_properties.subject = (
        f"Operational board pack for {branding.mine_name}"
    )
    presentation.core_properties.author = "Mine Manager AI"
    presentation.core_properties.company = branding.company_name
    presentation.core_properties.comments = (
        "Generated by Mine Manager AI from configured operational data."
    )
    presentation.core_properties.created = datetime.now()
    presentation.core_properties.modified = datetime.now()

    _build_cover_slide(
        presentation,
        branding,
    )

    _build_kpi_summary_slide(
        presentation,
        branding,
        production_rows,
        fleet_rows,
        plant_rows,
        safety_rows,
    )

    _build_production_slide(
        presentation,
        branding,
        production_rows,
    )

    _build_operations_overview_slide(
        presentation,
        branding,
        fleet_rows,
        plant_rows,
        safety_rows,
    )

    _build_risk_slide(
        presentation,
        branding,
        production_rows,
        fleet_rows,
        plant_rows,
        safety_rows,
    )

    _build_actions_slide(
        presentation,
        branding,
        executive_actions,
    )

    _build_recommendation_slide(
        presentation,
        branding,
        production_rows,
        fleet_rows,
        plant_rows,
        safety_rows,
    )

    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)

    return buffer
